#!/usr/bin/env python3
"""
Build the AI Kavach / Terrier Cyber Quest 2026 five-slide submission for Brahmadatta AI.

Editable source for `brahmadatta-ai-aikavach-2026.pptx`. Mirrors the HTML deck
(`brahmadatta-ai-aikavach-2026.html`) and the Command Center visual system
(docs/09-company/04-design-system.md, packages/ui-components/tokens.css).

    pip install python-pptx
    python submission/build_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.oxml.ns import qn
import copy

# ----------------------------------------------------------------- palette ---
FIELD      = RGBColor(0x14, 0x1C, 0x8C)
WHITE      = RGBColor(0xF4, 0xF3, 0xEE)
VEIL       = RGBColor(0x9F, 0xA1, 0xC9)
RULE       = RGBColor(0x72, 0x76, 0xB5)
RULE_FAINT = RGBColor(0x4A, 0x50, 0xA4)
VERIFIED   = RGBColor(0x39, 0xE0, 0x8A)
WARNING    = RGBColor(0xFF, 0xB0, 0x20)
CRITICAL   = RGBColor(0xFF, 0x6B, 0x66)

SERIF = "Georgia"          # Didone-ish, universally available fallback for Instrument Serif
MONO  = "Consolas"         # utility face; PowerPoint substitutes a mono if absent

EMU_IN = 914400
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


# ------------------------------------------------------------- primitives ---
def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = FIELD
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def _no_shadow(shp):
    try:
        shp.shadow.inherit = False
    except Exception:
        pass


def rect(s, x, y, w, h, line=RULE, weight=0.75, fill=None, dash=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(weight)
        if dash:
            d = sp.line._get_or_add_ln()
            pd = d.makeelement(qn('a:prstDash'), {'val': dash})
            d.append(pd)
    _no_shadow(sp)
    return sp


def hline(s, x, y, w, color=RULE, weight=0.75):
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    _no_shadow(ln)
    return ln


def vline(s, x, y, h, color=RULE, weight=0.75):
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x), Inches(y + h))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    _no_shadow(ln)
    return ln


def arrow(s, x1, y1, x2, y2, color=RULE, weight=0.75):
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    lnEl = ln.line._get_or_add_ln()
    tail = lnEl.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    lnEl.append(tail)
    _no_shadow(ln)
    return ln


def frame(s):
    """Page frame inset 0.34in, with white corner crop ticks over a broken corner."""
    inset = 0.34
    rect(s, inset, inset, SW - 2 * inset, SH - 2 * inset, line=RULE, weight=0.75)
    g = 0.09          # gap: field square that breaks the corner
    for cx, cy in [(inset, inset), (SW - inset, inset), (inset, SH - inset), (SW - inset, SH - inset)]:
        rect(s, cx - g, cy - g, 2 * g, 2 * g, line=None, fill=FIELD)
    t = 0.16          # tick length
    for cx, cy, sx, sy in [(inset, inset, 1, 1), (SW - inset, inset, -1, 1),
                           (inset, SH - inset, 1, -1), (SW - inset, SH - inset, -1, -1)]:
        hline(s, cx if sx > 0 else cx - t, cy, t, color=WHITE, weight=1.0)
        vline(s, cx, cy if sy > 0 else cy - t, t, color=WHITE, weight=1.0)


def tb(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def _apply(run, font, size, color, *, bold=False, italic=False, spacing=None, caps=False):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.language_id = MSO_LANGUAGE_ID.ENGLISH_UK
    if spacing is not None:
        rPr = run._r.get_or_add_rPr()
        rPr.set('spc', str(int(spacing)))
    if caps:
        rPr = run._r.get_or_add_rPr()
        rPr.set('cap', 'all')


def para(tf, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    return p


def line_runs(tf, segments, *, size=11, font=MONO, color=WHITE, space_after=4,
              space_before=0, leading=1.12, align=PP_ALIGN.LEFT, first=False,
              spacing=None, caps=False, bullet=False):
    """segments: list of (text, color|None, {overrides}) or plain str."""
    p = para(tf, first=first)
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.line_spacing = leading
    if bullet:
        run = p.add_run(); run.text = "·  "
        _apply(run, font, size, RULE)
    for seg in segments:
        if isinstance(seg, str):
            txt, col, ov = seg, color, {}
        else:
            txt, col, ov = (list(seg) + [None, {}])[:3]
            col = col or color
            ov = ov or {}
        run = p.add_run()
        run.text = txt
        _apply(run, ov.get("font", font), ov.get("size", size), col,
               bold=ov.get("bold", False), italic=ov.get("italic", False),
               spacing=ov.get("spacing", spacing), caps=ov.get("caps", caps))
    return p


def eyebrow(s, left, right):
    tf = tb(s, 0.72, 0.62, SW - 1.44, 0.34)
    p = tf.paragraphs[0]
    p.line_spacing = 1.0
    r = p.add_run(); r.text = left
    _apply(r, MONO, 9.5, VEIL, spacing=140, caps=True)
    tf2 = tb(s, SW - 5.2 - 0.72, 0.62, 5.2, 0.34)
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run(); r2.text = right
    _apply(r2, MONO, 9.5, WHITE, spacing=140, caps=True)


def title(s, y, segments, size=44, width=None):
    tf = tb(s, 0.72, y, width or (SW - 1.9), 1.9)
    p = tf.paragraphs[0]
    p.line_spacing = 0.98
    for txt, col in segments:
        r = p.add_run(); r.text = txt
        _apply(r, SERIF, size, col)
    return tf


def crop_panel(s, x, y, w, h, tint=RULE):
    t = 0.12
    for cx, cy, sx, sy in [(x, y, 1, 1), (x + w, y, -1, 1), (x, y + h, 1, -1), (x + w, y + h, -1, -1)]:
        hline(s, cx if sx > 0 else cx - t, cy, t, color=tint, weight=0.9)
        vline(s, cx, cy if sy > 0 else cy - t, t, color=tint, weight=0.9)


def label(tf, text, *, first=False, space_after=6):
    line_runs(tf, [(text, VEIL, {"spacing": 120})], size=9.5, space_after=space_after,
              caps=True, first=first)


# ================================================================ SLIDE 1 ===
def slide1():
    s = slide(); frame(s)
    eyebrow(s, "[ Slide 01 / 05 ]   Introduction · Ideation · Brief Description", "[ Brahmadatta AI ]")

    title(s, 1.18, [("Autonomous armor for software.", WHITE)], size=46, width=8.6)

    tf = tb(s, 0.72, 2.72, 8.7, 1.5)
    line_runs(tf, [
        ("An ", VEIL), ("authorised, defensive Cyber-Reasoning System", WHITE),
        (" for AI Kavach. Given a C/C++ repository you own, Brahmadatta finds a real "
         "memory-safety defect with its own fuzzing, drafts a patch with a ", VEIL),
        ("self-hosted", WHITE),
        (" model, and lets deterministic tools — not the model — deliver the verdict.", VEIL),
    ], size=12.5, leading=1.4, first=True)

    # three panels
    cols = [
        ("[ Problem ]", [
            [("Scanners and “AI security” tools stop at the report. They hand a human a "
              "finding and a suggested patch.", WHITE)],
            [("A suggested patch, ", WHITE), ("accepted on model confidence", WHITE, {"bold": False}),
             (", is exactly how a plausible-looking wrong fix ships to production.", WHITE)],
        ]),
        ("[ Idea & Motivation ]", [
            [("Close the loop, and make the machine check its own work.", WHITE)],
            [("Discover → minimise to a deterministic reproducer → patch with a local "
              "model → ", WHITE),
             ("let the reproducer, the regression suite, static checks and renewed fuzzing decide.",
              WHITE)],
            [("Nothing escalates on a guess.", VEIL)],
        ]),
        ("[ Brahmadatta AI, in one line ]", [
            [("authorize → ingest → baseline → analyze → correlate → "
              "stress-test → patch → verify → export evidence", VEIL, {"size": 9.5})],
            [("One autonomous pipeline. Three evidence-gated tiers. One operator-visible ", WHITE),
             ("Command Center", WHITE), (".", WHITE)],
        ]),
    ]
    px, pw, gap = 0.72, 3.86, 0.30
    py, ph = 4.30, 1.72
    for i, (lbl, body) in enumerate(cols):
        x = px + i * (pw + gap)
        crop_panel(s, x, py, pw, ph)
        tf = tb(s, x + 0.18, py + 0.16, pw - 0.36, ph - 0.3)
        label(tf, lbl, first=True)
        for j, segs in enumerate(body):
            line_runs(tf, segs, size=9.7, leading=1.32, space_after=4)

    hline(s, 0.72, 6.28, SW - 1.44, color=RULE)

    b2 = [
        ("[ The Differentiator ]", [
            ("The same pipeline visibly ", VEIL), ("rejects", CRITICAL),
            (" a patch that eliminates the crash but breaks the program's behaviour. Every "
             "competitor shows a patch that worked; we show our system refusing one that only "
             "looked like it did.", VEIL)]),
        ("[ Defensive / Authorised-Use Boundary ]", [
            ("Authorised repositories and isolated environments only. Written authorization record "
             "+ server-verified snapshot hash. ", VEIL),
            ("No public-target scanning, no exploit deployment, no automatic production merge.", WHITE),
            (" Sandbox runs --network none; repository content never leaves for a hosted "
             "inference API.", VEIL)]),
    ]
    for i, (lbl, segs) in enumerate(b2):
        x = 0.72 + i * (6.0 + 0.4)
        tf = tb(s, x, 6.44, 6.0, 0.95)
        label(tf, lbl, first=True, space_after=4)
        line_runs(tf, segs, size=9.6, leading=1.32)
    return s


# ================================================================ SLIDE 2 ===
def slide2():
    s = slide(); frame(s)
    eyebrow(s, "[ Slide 02 / 05 ]   Detailed Methodology", "[ The Mission — Nine Steps ]")
    title(s, 1.1, [("One pipeline. ", WHITE), ("Unattended after a single ", VEIL),
                   ("[ START ].", VEIL)], size=40, width=11.6)
    hline(s, 0.72, 2.28, SW - 1.44, color=RULE)

    cols = [
        ("[ Tier 1 · Deterministic Triage — CPU only, no model ]", [
            [("01 Authorize + snapshot.", WHITE), (" Immutable snapshot; sha256 recomputed "
             "server-side and matched.", VEIL)],
            [("02 Baseline.", WHITE), (" configure → build → CTest on the pristine tree. "
             "Pass/fail counts recorded — the denominator every later verdict is checked "
             "against.", VEIL)],
            [("03 Analyze.", WHITE), (" Semgrep (offline vendored C/C++ ruleset) + compiler "
             "warnings parsed into structured findings; git-history summary.", VEIL)],
        ]),
        ("[ Tier 2 · Destructive Sandbox + Lightweight Patch ]", [
            [("04 Stress-test.", WHITE), (" ASan/UBSan build + libFuzzer campaign inside a "
             "locked-down container. Sanitizer-confirmed heap-buffer-overflow with a stack "
             "trace; crashes deduplicated.", VEIL)],
            [("05 Correlate + minimise.", WHITE), (" Bind the crash to a source location; minimise "
             "the input to a reproducer that replays 5/5 from a clean build; persist it durably.",
             VEIL)],
            [("06 Patch.", WHITE), (" A self-hosted CodeLlama 7B model gets the crash report + "
             "localised source. Patch policy (single file, allowlist, changed-line cap) must pass "
             "before the candidate is compiled.", VEIL)],
        ]),
        ("[ Verification & Proof-of-Fix Loop ]", [
            [("07 Verify — every policy-passing candidate, identical gate sequence.", WHITE),
             (" The verifier is provenance-blind: it cannot see the model, its confidence, or its "
              "rationale.", VEIL)],
            [("08 Proof-of-fix.", WHITE), (" The verdict is derived from the gate matrix, never "
             "from confidence.", VEIL)],
            [("09 Export + teardown.", WHITE), (" report.md / .json, gate matrix, content-addressed "
             "artifacts, sha256 manifest. Sandbox + model-host release confirmed, zero strays.",
             VEIL)],
        ]),
    ]
    px, pw, gap = 0.72, 3.86, 0.30
    py, ph = 2.5, 2.62
    for i, (lbl, body) in enumerate(cols):
        x = px + i * (pw + gap)
        crop_panel(s, x, py, pw, ph)
        tf = tb(s, x + 0.18, py + 0.16, pw - 0.36, ph - 0.3)
        label(tf, lbl, first=True)
        for segs in body:
            line_runs(tf, segs, size=9.2, leading=1.28, space_after=4, bullet=True)

    hline(s, 0.72, 5.36, SW - 1.44, color=WHITE, weight=1.6)

    # gate matrix + explanation
    tf = tb(s, 0.72, 5.54, 6.0, 1.6)
    label(tf, "[ Step 08 · Two candidates, one gate matrix, one run ]", first=True, space_after=5)
    for glyph, gc, name, meta in [
        ("+", VERIFIED, "COMPILE", "cmake · exit 0"),
        ("+", VERIFIED, "REPRODUCER_ELIMINATED", "replay clean"),
        ("+", VERIFIED, "REGRESSION_PRESERVED  (patch A)", "8/8 ctest"),
        ("×", CRITICAL, "REGRESSION_PRESERVED  (patch B)", "1 of 8 failed"),
        ("—", WARNING, "STATIC_DELTA · RENEWED_FUZZING", "disclosed, with reason"),
    ]:
        line_runs(tf, [(glyph + "  ", gc), (name, WHITE), ("   " + meta, VEIL, {"size": 8})],
                  size=9.5, leading=1.5, space_after=1)

    tf2 = tb(s, 7.1, 5.54, 5.5, 1.7)
    label(tf2, "[ What the loop catches ]", first=True, space_after=5)
    line_runs(tf2, [
        ("Patch A — the correct fix — is ", VEIL), ("VERIFIED", VERIFIED),
        (". Patch B kills the crash by deleting the feature that was overflowing: it passes "
         "compile, eliminates the reproducer, then ", VEIL),
        ("fails a regression test", CRITICAL),
        (" (six assertions in the tab-expansion case) → ", VEIL),
        ("REJECTED", CRITICAL),
        (", beside the accepted one. An overfit patch that fools every static gate is caught by ",
         VEIL),
        ("renewed fuzzing", WHITE), (" of the patched build.", VEIL),
    ], size=9.4, leading=1.34)

    tf3 = tb(s, 0.72, 7.02, SW - 1.44, 0.3)
    line_runs(tf3, [("[ Nothing escalates a tier on a guess — work moves up only when the "
                     "evidence in hand justifies the cost ]", VEIL, {"spacing": 60})],
              size=9, caps=True, first=True)
    return s


# ================================================================ SLIDE 3 ===
def dbox(s, x, y, w, h, title_txt, lines, fin=False):
    rect(s, x, y, w, h, line=(WARNING if fin else RULE),
         weight=0.75, dash=('dash' if fin else None))
    tf = tb(s, x + 0.12, y + 0.09, w - 0.24, h - 0.18)
    line_runs(tf, [(title_txt, (WARNING if fin else WHITE))], size=8.6, space_after=2,
              leading=1.1, first=True)
    for ln in lines:
        line_runs(tf, [(ln, VEIL)], size=7.2, leading=1.16, space_after=1)


def slide3():
    s = slide(); frame(s)
    eyebrow(s, "[ Slide 03 / 05 ]   Technology Stack · Architecture · Flow", "[ As Built ]")
    title(s, 1.08, [("Four processes. ", WHITE), ("No broker. No hosted API.", VEIL)],
          size=38, width=11.6)
    hline(s, 0.72, 2.2, SW - 1.44, color=RULE)

    # ---- diagram (left) ----
    ox, oy = 0.72, 2.42
    dbox(s, ox, oy, 3.0, 0.62, "operator browser", ["1440×900+ · Command Center"])
    dbox(s, ox, oy + 0.95, 3.0, 0.82, "nginx :443 / :8080",
         ["TLS · serves Astro build", "proxy_buffering off on /events"])
    arrow(s, ox + 1.5, oy + 0.62, ox + 1.5, oy + 0.95)

    dbox(s, ox, oy + 2.05, 4.35, 1.06, "control-api — Django 5.2 · django-ninja",
         ["uvicorn ASGI · generated OpenAPI contract in CI",
          "SSE fan-out · validates operator commands",
          "READS everything · RUNS nothing"])
    arrow(s, ox + 1.5, oy + 1.77, ox + 1.5, oy + 2.05)

    dbox(s, ox, oy + 3.35, 4.35, 0.72, "PostgreSQL 16",
         ["missions · events · jobs · evidence"])
    arrow(s, ox + 2.1, oy + 3.11, ox + 2.1, oy + 3.35)

    dbox(s, ox, oy + 4.32, 5.4, 1.12, "orchestrator — run_orchestrator",
         ["sole writer of Mission.state · 18-state machine",
          "tick loop: transition · enqueue · reap · watchdog",
          "queue = Postgres SELECT … FOR UPDATE SKIP LOCKED"])
    arrow(s, ox + 1.1, oy + 4.07, ox + 1.1, oy + 4.32)
    arrow(s, ox + 2.7, oy + 4.32, ox + 2.7, oy + 4.07)

    wx = ox + 5.75
    dbox(s, wx, oy + 2.05, 4.6, 0.92, "worker — JobKind dispatch",
         ["BASELINE · ANALYZE · SANITIZER_BUILD · FUZZ · MINIMIZE",
          "CORRELATE · PATCH_GENERATE · VERIFY · EXPORT · TEARDOWN"])
    dbox(s, wx, oy + 3.2, 4.6, 1.28, "sandbox — packages/sandbox",
         ["subprocess Jail  +  ContainerJail",
          "--network none · --cap-drop ALL · non-root",
          "read-only rootfs · no docker socket",
          "CMake/CTest · ASan/UBSan · libFuzzer"])
    dbox(s, wx, oy + 4.68, 4.6, 0.86, "model-host — CodeLlama 7B (Ollama)",
         ["loopback / internal-only · bearer-token sidecar",
          "gateway = the only inference client in the tree"])
    # worker <-> orchestrator (job claim)
    arrow(s, ox + 5.4, oy + 4.88, wx, oy + 4.88)
    hline(s, wx - 0.35, oy + 4.88, 0.35, color=RULE)
    vline(s, wx - 0.35, oy + 2.97, oy + 4.88 - (oy + 2.97), color=RULE)
    arrow(s, wx + 2.3, oy + 2.97, wx + 2.3, oy + 3.2)
    arrow(s, wx + 2.3, oy + 4.48, wx + 2.3, oy + 4.68)

    dbox(s, ox, oy + 5.66, 5.4, 0.62, "content-addressed artifact store",
         ["sha256 · mode 0600 · encrypted volume"])

    tf = tb(s, wx, oy - 0.02, 4.6, 0.24)
    line_runs(tf, [("internal: true — no route to the internet", VEIL)], size=7.4,
              align=PP_ALIGN.RIGHT, first=True)
    vline(s, wx - 0.16, oy + 1.9, 4.0, color=RULE_FAINT)

    # ---- stack legend (right) ----
    lx = ox + 10.55
    tf = tb(s, lx, oy + 0.0, 1.9, 0.24)  # placeholder to keep spacing
    # Actually the diagram is wide; put legend below title full width? Keep compact at far right.
    # ---- legend as a strip under the diagram would overflow; instead a right rail ----
    lx = 10.75
    tf = tb(s, lx, 2.42, 2.45, 4.4)
    label(tf, "[ Stack — implemented ]", first=True, space_after=4)
    for name in ["Astro 7 + React islands", "Django · django-ninja · Pydantic 2",
                 "PostgreSQL 16 · ORM + migrations", "18-state machine · DB queue",
                 "Semgrep · compiler warnings", "libFuzzer · ASan/UBSan · dedup",
                 "rootless Docker · Jail + ContainerJail", "self-hosted CodeLlama 7B",
                 "5-gate matrix · no confidence arg"]:
        line_runs(tf, [("+ ", VERIFIED), (name, WHITE)], size=7.6, leading=1.24, space_after=2)
    label(tf, "[ Finale / roadmap ]", space_after=4)
    for name in ["Live-model patch in the full run", "Semgrep + renewed-fuzz in a live mission",
                 "ContainerJail default everywhere", "Tier 3 — designed path only · unnamed"]:
        line_runs(tf, [("› ", WARNING), (name, VEIL)], size=7.6, leading=1.24, space_after=2)

    tf = tb(s, 0.72, 7.04, SW - 1.44, 0.3)
    line_runs(tf, [("Fifteen deployable units in the original design collapse to four application "
                    "processes. The worker is the only thing that holds repository content and the "
                    "only caller of the model.", VEIL)], size=8, first=True)
    return s


# ================================================================ SLIDE 4 ===
def slide4():
    s = slide(); frame(s)
    eyebrow(s, "[ Slide 04 / 05 ]   Salient Features & Novelty",
            "[ Resource Utilisation · Novelty · Lightweight ]")
    title(s, 1.05, [("The system rejects its own patch — ", WHITE),
                    ("and the rejection is a real regression-test failure, not a scripted demo "
                     "beat.", VEIL)], size=30, width=11.9)
    hline(s, 0.72, 2.5, SW - 1.44, color=RULE)

    cols = [
        ("[ Novelty — verification that overrules the model ]", [
            [("Evidence-gated escalation.", WHITE), (" Nothing reaches the model until "
             "deterministic tiers produce a confirmed, minimised finding.", VEIL)],
            [("No confidence path.", WHITE), (" derive_verdict() takes one argument — the gate "
             "matrix. A VERIFIED record over a failing gate cannot be constructed.", VEIL)],
            [("Provenance-blind verifier.", WHITE), (" The operator-supplied “bad” "
             "candidate takes a genuinely identical path.", VEIL)],
            [("Overfit caught by renewed fuzzing", WHITE), (" of the patched build — not by a "
             "human already knowing which input to try.", VEIL)],
            [("Disclosure as a feature.", WHITE), (" A gate that did not run is as loud as a "
             "failure; the verdict carries its denominator.", VEIL)],
        ]),
        ("[ Posture — self-hosted · authorised · isolated ]", [
            [("Repository content never reaches a hosted API.", WHITE), (" Enforced three ways: "
             "internal-only compose network, a single-inference-client source test, a boot-time "
             "system check that refuses a hosted endpoint.", VEIL)],
            [("Rootless isolated execution", WHITE), (" — --network none, all capabilities "
             "dropped, non-root, read-only rootfs. Egress denial proven by a live DNS+TCP test.",
             VEIL)],
            [("Authorised repositories only", WHITE), (" — write-once authorization record, "
             "server-verified snapshot hash. No public scanning, no exploit deployment, no "
             "auto-merge.", VEIL)],
            [("Operator-visible Command Center", WHITE), (" — every panel reads the same "
             "stream that produces the evidence bundle.", VEIL)],
        ]),
        ("[ Resource utilisation & lightweight design ]", [
            [("CPU-first.", WHITE), (" Rented GPU cut entirely. A local 7B model with a "
             "process-level lease and confirmed teardown.", VEIL)],
            [("4 application processes", WHITE), (" vs 15 in the original design. Postgres SKIP "
             "LOCKED instead of a message broker. No Redis, no S3, offline Semgrep ruleset.", VEIL)],
            [("Runs on one machine, offline.", WHITE), (" Orchestrator / worker / gateway / "
             "evidence are modules in one Django project.", VEIL)],
        ]),
    ]
    px, pw, gap = 0.72, 3.86, 0.30
    py, ph = 2.72, 3.3
    for i, (lbl, body) in enumerate(cols):
        x = px + i * (pw + gap)
        crop_panel(s, x, py, pw, ph)
        tf = tb(s, x + 0.18, py + 0.16, pw - 0.36, ph - 0.3)
        label(tf, lbl, first=True)
        for segs in body:
            line_runs(tf, segs, size=8.6, leading=1.24, space_after=4, bullet=True)

    # stat row inside col 3
    tf = tb(s, px + 2 * (pw + gap) + 0.18, py + 2.5, pw - 0.36, 0.72)
    p = tf.paragraphs[0]; p.line_spacing = 1.0
    for txt, col, sz in [("~48s", WHITE, 22), ("  full pipeline   ", VEIL, 8),
                         ("0.3s", WHITE, 22), ("  fuzz to seeded defect", VEIL, 8)]:
        r = p.add_run(); r.text = txt
        _apply(r, SERIF if sz > 12 else MONO, sz, col)

    hline(s, 0.72, 6.28, SW - 1.44, color=WHITE, weight=1.6)
    tf = tb(s, 0.72, 6.44, SW - 1.44, 0.4)
    line_runs(tf, [("[ Not another “LLM + scanner” wrapper — a verification layer "
                    "that is demonstrated overruling the model ]", VEIL, {"spacing": 60})],
              size=9, caps=True, first=True)
    return s


# ================================================================ SLIDE 5 ===
def slide5():
    s = slide(); frame(s)
    eyebrow(s, "[ Slide 05 / 05 ]   Final Deliverables", "[ Evidence Run 6 · 2026-08-20 ]")
    title(s, 1.08, [("Demonstrated today. ", WHITE), ("Sharpened at the finale.", VEIL)],
          size=40, width=11.6)
    hline(s, 0.72, 2.24, SW - 1.44, color=RULE)

    # left panel — done
    crop_panel(s, 0.72, 2.46, 6.2, 3.7, tint=VERIFIED)
    tf = tb(s, 0.92, 2.64, 5.85, 3.4)
    label(tf, "[ What the prototype demonstrably does — live, unattended, real HTTP API ]",
          first=True, space_after=6)
    for segs in [
        [("Full nine-step pipeline end-to-end in ", VEIL), ("47.75 s", WHITE), (" wall-clock.", VEIL)],
        [("A self-discovered, ASan-confirmed heap-buffer-overflow in an authorised C target, with "
          "a durable, deterministic reproducer.", VEIL)],
        [("Two verdicts from one mission", WHITE), (" — one patch → ", VEIL),
         ("VERIFIED", VERIFIED), (" (first end-to-end in project history); one plausible crash-only "
         "patch → ", VEIL), ("REJECTED", CRITICAL), (" on a real regression failure, same "
         "gate matrix.", VEIL)],
        [("An exported evidence bundle", WHITE), (" — snapshot hash, crash report, minimised "
         "input, both diffs, both gate matrices, both verdicts — sha256-manifested and "
         "independently read back by someone who did not build the system.", VEIL)],
        [("Confirmed teardown, zero stray containers.", VEIL)],
        [("Self-hosted CodeLlama 7B: ", VEIL), ("10 / 10", WHITE), (" generation attempts returned "
         "schema- and policy-valid candidates.", VEIL)],
        [("Command Center rebuilt to the approved spec; ≈900+ automated tests; CI on a real "
          "PostgreSQL; egress denial proven live.", VEIL)],
    ]:
        line_runs(tf, segs, size=9.0, leading=1.26, space_after=4, bullet=True)

    # right column — output + finale
    crop_panel(s, 7.2, 2.46, 5.4, 1.35)
    tf = tb(s, 7.4, 2.62, 5.05, 1.1)
    label(tf, "[ Expected competition output ]", first=True, space_after=4)
    line_runs(tf, [("A reproducible C/C++ pipeline that, on an authorised repository, produces a "
                    "confirmed finding, a minimal patch, a deterministic verdict, and a ", VEIL),
                   ("portable evidence bundle a judge can audit offline", WHITE), (".", VEIL)],
              size=8.6, leading=1.24)

    crop_panel(s, 7.2, 3.98, 5.4, 2.18, tint=WARNING)
    tf = tb(s, 7.4, 4.14, 5.05, 1.95)
    label(tf, "[ Refined during the 36-hour Grand Finale ]", first=True, space_after=4)
    for ln in ["Live-model patch generation inside the full unattended run (finale hardware clears "
               "the dev-VM memory ceiling).",
               "Semgrep + renewed-fuzzing gates exercised inside a full live mission.",
               "ContainerJail as the default isolation for every stage, every deployment.",
               "Three timed rehearsals with failure injection: GPU-unavailable, target won't "
               "build, stage hangs.",
               "A real open-source C target with a known historical CVE, alongside the "
               "purpose-built one.",
               "Human-recorded fallback demonstration."]:
        line_runs(tf, [(ln, VEIL)], size=8.0, leading=1.18, space_after=3, bullet=True)

    hline(s, 0.72, 6.34, SW - 1.44, color=WHITE, weight=1.6)
    tf = tb(s, 0.72, 6.5, SW - 1.44, 0.8)
    line_runs(tf, [
        ("PERFORMANCE ", VERIFIED), ("~48 s pipeline    ", VEIL),
        ("SPEED ", VERIFIED), ("sub-second reproducer replay    ", VEIL),
        ("PRECISION ", VERIFIED), ("deterministic gates + demonstrated rejection", VEIL),
    ], size=8.8, leading=1.3, first=True)
    line_runs(tf, [
        ("FUNCTIONALITY ", VERIFIED), ("nine-step loop + live Command Center    ", VEIL),
        ("SCALABILITY ", VERIFIED), ("stateless processes · SKIP LOCKED queue · adapter "
         "generalised beyond the demo target", VEIL),
    ], size=8.8, leading=1.3)
    return s


slide1(); slide2(); slide3(); slide4(); slide5()

out = __file__.rsplit("/", 1)[0] + "/brahmadatta-ai-aikavach-2026.pptx"
prs.save(out)
print("wrote", out)
